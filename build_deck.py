#!/usr/bin/env python3
"""
Contabo Branded PPTX Builder — Standalone executable.
Downloads assets from GitHub, builds deck from JSON content file.

USAGE:
    python3 build_deck.py content.json output.pptx

content.json format:
{
  "cover_topic": null,   // or: "growth","privacy","search","security","server","database"
  "slides": [
    {"type": "cover", "title": "...", "subtitle": "...", "meta": "..."},
    {"type": "2col", "title": "...", "subtitle": "...",
     "left_label": "...", "left_body": "...",
     "right_label": "...", "right_body": "...", "footer": "..."},
    {"type": "3col", "title": "...", "subtitle": "...",
     "cols": [["LABEL","subhead","body"], ...], "footer": "..."},
    {"type": "content", "title": "...", "subtitle": "...", "body": "..."},
    {"type": "divider", "title": "...", "subtitle": "..."},
    {"type": "closing"}
  ]
}
"""

import sys, json, os, tempfile, urllib.request, base64
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ══════════════════════════════════════════════════════════════════════════════
# CONSTANTS
# ══════════════════════════════════════════════════════════════════════════════

SC = 33.8667 / 1920
def px(v): return Cm(round(v * SC, 5))

W = Cm(33.8667)
H = Cm(19.05)

C_GRAY_2  = RGBColor(0x19, 0x23, 0x2D)
C_GRAY_1  = RGBColor(0x32, 0x3C, 0x46)
C_GRAY    = RGBColor(0x5A, 0x64, 0x73)
C_GRAY_P1 = RGBColor(0x82, 0x8C, 0x9B)
C_GRAY_P4 = RGBColor(0xF5, 0xF5, 0xFA)
C_BLUE_2  = RGBColor(0x0A, 0x50, 0x96)
C_BLUE    = RGBColor(0x00, 0xAA, 0xEB)
C_ORANGE  = RGBColor(0xFF, 0xB4, 0x00)
C_GREEN   = RGBColor(0x14, 0xCD, 0x82)
C_RED     = RGBColor(0xD7, 0x19, 0x19)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

GITHUB_BASE = "https://raw.githubusercontent.com/markosansa/ppt/main/"

GITHUB_ASSETS = {
    "LOGO_WHITE":      "logo-white.png",
    "LOGO_PRIMARY":    "logo-primary.png",
    "LOGO_ICON_WHT":   "logo-icon-white-transparent.png",
    "LOGO_ICON_CLR":   "logo-icon-color-transparent.png",
    "COVER_1":         "Master-cover-1.png",
    "COVER_2":         "Master-cover-2.png",
    "COVER_3":         "Master-cover-3.png",
    "CUSTOM_DATABASE": "Custom-Cover-DATABASE.png",
    "CUSTOM_GROWTH":   "Custom-Cover-GROWTH.png",
    "CUSTOM_PRIVACY":  "Custom-Cover-PRIVACY.png",
    "CUSTOM_SEARCH":   "Custom-Cover-SEARCH.png",
    "CUSTOM_SECURITY": "Custom-Cover-SECURITY.png",
    "CUSTOM_SERVER":   "Custom-Cover-SERVER.png",
}

COVER_MAP = {
    "growth":   "CUSTOM_GROWTH",
    "privacy":  "CUSTOM_PRIVACY",
    "search":   "CUSTOM_SEARCH",
    "security": "CUSTOM_SECURITY",
    "server":   "CUSTOM_SERVER",
    "database": "CUSTOM_DATABASE",
}

# ══════════════════════════════════════════════════════════════════════════════
# ASSET LOADING
# ══════════════════════════════════════════════════════════════════════════════

def download_assets():
    """Download all brand assets from GitHub. Returns dict of key -> local path."""
    paths = {}
    for key, filename in GITHUB_ASSETS.items():
        url = GITHUB_BASE + filename
        ext = filename.rsplit(".", 1)[-1]
        tmp = tempfile.NamedTemporaryFile(suffix="." + ext, delete=False)
        try:
            urllib.request.urlretrieve(url, tmp.name)
            if os.path.getsize(tmp.name) > 512:
                paths[key] = tmp.name
            else:
                print(f"  WARNING: {key} too small ({os.path.getsize(tmp.name)}B)", file=sys.stderr)
        except Exception as e:
            print(f"  WARNING: {key} download failed: {e}", file=sys.stderr)
    return paths


def load_master(master_pptx_path=None):
    """
    Open the master .pptx template.
    If master_pptx_path is provided, use it.
    Otherwise, download from GitHub.
    Returns a Presentation with all content slides removed.
    """
    if master_pptx_path and os.path.exists(master_pptx_path):
        prs = Presentation(master_pptx_path)
    else:
        # Try downloading from GitHub
        url = GITHUB_BASE + "Template-master-SLIM.pptx"
        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        try:
            urllib.request.urlretrieve(url, tmp.name)
            prs = Presentation(tmp.name)
        except Exception:
            raise RuntimeError(
                "FATAL: Cannot load master template. Upload Template-master.pptx "
                "or ensure Template-master-SLIM.pptx is on GitHub."
            )

    # Remove any existing content slides
    xml_slides = prs.slides._sldIdLst
    while len(xml_slides) > 0:
        rId = xml_slides[0].get(qn("r:id"))
        if rId:
            prs.part.drop_rel(rId)
        del xml_slides[0]

    # Verify this is the Contabo master
    layout_names = [l.name for l in prs.slide_masters[0].slide_layouts]
    if len(layout_names) < 20 or "Cover" not in layout_names:
        raise RuntimeError(
            f"FATAL: Wrong master template loaded! Found {len(layout_names)} layouts. "
            f"Expected 27+ Contabo layouts including 'Cover', '2_Cover', 'Closing', etc."
        )

    return prs


# ══════════════════════════════════════════════════════════════════════════════
# SHAPE HELPERS — every fix baked in, never modified by Claude
# ══════════════════════════════════════════════════════════════════════════════

def _strip_effects(el):
    """
    Remove ALL shadow/effect sources from a shape element.

    CRITICAL: shapes created via add_shape() get an auto-generated <p:style>
    block containing <a:effectRef idx="2"> which references the THEME's
    shadow style. Stripping effectLst from spPr alone does NOT remove this.
    The entire <p:style> must be removed, plus an explicit empty <a:effectLst/>
    must be added to spPr to prevent inheritance.
    """
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"

    # 1. Remove explicit effect elements in spPr
    for ns in [ns_a, ns_p]:
        for tag in ["effectLst", "effectDag", "outerShdw", "innerShdw",
                     "glow", "reflection", "softEdge"]:
            for e in el.findall(f".//{{{ns}}}{tag}"):
                parent = e.getparent()
                if parent is not None:
                    parent.remove(e)

    # 2. Remove <p:style> entirely — THE fix for theme-inherited shadows
    for style in el.findall(f"{{{ns_p}}}style"):
        el.remove(style)

    # 3. Ensure spPr has explicit empty effectLst
    spPr = el.find(f"{{{ns_p}}}spPr")
    if spPr is not None and spPr.find(f"{{{ns_a}}}effectLst") is None:
        etree.SubElement(spPr, f"{{{ns_a}}}effectLst")


def add_img(slide, path, l, t, w, h):
    """Add image with all effects stripped."""
    pic = slide.shapes.add_picture(path, l, t, w, h)
    _strip_effects(pic._element)
    pic.line.fill.background()
    return pic


def bg_img(slide, path):
    """Full-bleed background image — sent to back of z-order."""
    p = slide.shapes.add_picture(path, 0, 0, W, H)
    _strip_effects(p._element)
    slide.shapes._spTree.remove(p._element)
    slide.shapes._spTree.insert(2, p._element)
    return p


def box(slide, l, t, w, h, fill=None):
    """Rectangle shape with optional solid fill. No shadows, no line."""
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background()
    _strip_effects(s._element)
    return s


def txt(slide, l, t, w, h, text, pt, bold=False, color=None, align=PP_ALIGN.LEFT):
    """
    Text box with Calibri font, word wrap, minimal internal padding.
    Padding: lIns=rIns=45720 (0.05cm), tIns=bIns=22860 (0.025cm).
    """
    if color is None:
        color = C_GRAY_1
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    bp = tf._txBody.find(qn("a:bodyPr"))
    if bp is not None:
        bp.set("lIns", "45720")
        bp.set("rIns", "45720")
        bp.set("tIns", "22860")
        bp.set("bIns", "22860")
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(pt)
    r.font.bold = bold
    r.font.color.rgb = color
    _strip_effects(tb._element)
    return tb


# ══════════════════════════════════════════════════════════════════════════════
# LAYOUT HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def estimate_lines(title, chars_per_line=12):
    """
    Word-aware line counter. Counts how many lines text will wrap to,
    respecting word boundaries (words never split mid-word).
    """
    words = title.split()
    lines = 1
    current = 0
    for word in words:
        wl = len(word)
        if current == 0:
            current = wl
        elif current + 1 + wl <= chars_per_line:
            current += 1 + wl
        else:
            lines += 1
            current = wl
    return max(1, lines)


def grad_bg(slide):
    """White -> light-gray gradient background for content slides."""
    f = slide.background.fill
    f.gradient()
    f.gradient_angle = 135
    f.gradient_stops[0].position = 0
    f.gradient_stops[0].color.rgb = WHITE
    f.gradient_stops[1].position = 1
    f.gradient_stops[1].color.rgb = C_GRAY_P4


def logo_white(slide, A):
    """White Contabo logo — top-left, for dark backgrounds."""
    add_img(slide, A["LOGO_WHITE"], px(134), px(93), px(237), px(47))


def logo_primary(slide, A):
    """Color Contabo logo — bottom-right, for light backgrounds."""
    add_img(slide, A["LOGO_PRIMARY"], px(1769), px(1029), px(123), px(25))


def watermark(slide, A):
    """Cloud icon watermark — bottom-right, for dark backgrounds. NOT on closing."""
    add_img(slide, A["LOGO_ICON_WHT"], px(1482), px(812), px(500), px(424))


def title_block(slide, title, subtitle=None):
    """
    Accent square + title + optional subtitle for content slides.
    Height auto-expands based on title line count.
    """
    n_lines = estimate_lines(title, chars_per_line=28)
    LINE_H = px(38)  # 28pt Calibri ≈ 38px line height
    title_h = LINE_H * n_lines
    accent_h = max(px(60), title_h)
    box(slide, 0, px(75), px(57), accent_h, fill=C_BLUE)
    txt(slide, px(72), px(75), W - px(144), title_h, title, 28, bold=True, color=C_GRAY_2)
    bot = px(75) + accent_h
    if subtitle:
        txt(slide, px(72), bot + px(6), W - px(144), px(30), subtitle, 15, color=C_GRAY_1)
        bot += px(6) + px(30)
    return bot


def col_h(slide, x, y, w, label):
    """Column header: thin blue bar + all-caps label."""
    box(slide, x, y, w, px(4), fill=C_BLUE)
    txt(slide, x, y + px(8), w, px(18), label, 9, bold=True, color=C_GRAY)


def closing_bar(slide):
    """5px orange bar at very bottom — CLOSING SLIDE ONLY."""
    box(slide, 0, H - px(5), W, px(5), fill=C_ORANGE)


def new_slide(prs, layout_name="Blank"):
    """
    Add a new slide. Removes all placeholder shapes to prevent
    'Click to add title/text' overlap with our drawn content.
    """
    layout = None
    for l in prs.slide_masters[0].slide_layouts:
        if l.name == layout_name:
            layout = l
            break
    if layout is None:
        for l in prs.slide_masters[0].slide_layouts:
            if l.name == "Blank":
                layout = l
                break
    if layout is None:
        layout = prs.slide_layouts[6]

    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    return slide


def two_col(slide, bot, l_lbl, l_body, r_lbl, r_body, footer=None):
    """Two-column layout below the title block."""
    gap = px(48)
    cy = bot + px(28)
    cw = (W - px(144) - gap) / 2
    lx = px(72)
    rx = lx + cw + gap
    bh = H - cy - px(28) - (px(40) if footer else 0)
    col_h(slide, lx, cy, cw, l_lbl)
    txt(slide, lx, cy + px(30), cw, bh, l_body, 13, color=C_GRAY_1)
    col_h(slide, rx, cy, cw, r_lbl)
    txt(slide, rx, cy + px(30), cw, bh, r_body, 13, color=C_GRAY_1)
    if footer:
        txt(slide, px(72), H - px(38), W - px(144), px(32), footer, 10, color=C_GRAY_P1)


def three_col(slide, bot, cols, footer=None):
    """Three-column layout below the title block."""
    gap = px(40)
    cy = bot + px(28)
    cw = (W - px(144) - 2 * gap) / 3
    bh = H - cy - px(28) - (px(40) if footer else 0)
    for i, col_data in enumerate(cols):
        lbl = col_data[0] if len(col_data) > 0 else ""
        sub = col_data[1] if len(col_data) > 1 else None
        body = col_data[2] if len(col_data) > 2 else ""
        cx = px(72) + i * (cw + gap)
        col_h(slide, cx, cy, cw, lbl)
        if sub:
            txt(slide, cx, cy + px(30), cw, px(48), sub, 14, bold=True, color=C_GRAY_2)
            txt(slide, cx, cy + px(84), cw, bh - px(54), body, 12, color=C_GRAY_1)
        else:
            txt(slide, cx, cy + px(30), cw, bh, body, 12, color=C_GRAY_1)
    if footer:
        txt(slide, px(72), H - px(38), W - px(144), px(32), footer, 10, color=C_GRAY_P1)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN BUILD FUNCTION
# ══════════════════════════════════════════════════════════════════════════════

def build_deck(content_json_path, output_path, master_pptx_path=None):
    """
    Build a complete Contabo-branded deck from a JSON content file.
    """
    with open(content_json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    slides_data = data.get("slides", [])
    cover_topic = data.get("cover_topic", None)

    # 1. Download assets
    A = download_assets()

    # 2. Open master
    prs = load_master(master_pptx_path)

    # 3. Determine cover background
    cover_key = COVER_MAP.get(cover_topic, "COVER_1")
    cover_bg = A.get(cover_key, A.get("COVER_1"))

    # 4. Build each slide
    for sd in slides_data:
        stype = sd.get("type", "content")

        if stype == "cover":
            sl = new_slide(prs, "Cover")
            bg_img(sl, cover_bg)
            logo_white(sl, A)
            watermark(sl, A)
            # Word-aware: 52pt bold, ~12 chars/line at 960px
            cover_title = sd.get("title", "Title")
            cov_lines = estimate_lines(cover_title, chars_per_line=12)
            COV_LINE_H = px(130)
            title_h = COV_LINE_H * cov_lines
            txt(sl, px(134), px(280), px(960), title_h,
                cover_title, 52, bold=True, color=WHITE)
            y = px(280) + title_h + px(36)
            if sd.get("subtitle"):
                txt(sl, px(134), y, px(940), px(60),
                    sd["subtitle"], 24, color=WHITE)
                y += px(72)
            if sd.get("meta"):
                txt(sl, px(134), y, px(940), px(40),
                    sd["meta"], 16, color=WHITE)

        elif stype == "divider":
            sl = new_slide(prs, "Blank")
            bg_img(sl, A.get("COVER_2", cover_bg))
            logo_white(sl, A)
            watermark(sl, A)
            div_title = sd.get("title", "Section")
            div_lines = estimate_lines(div_title, chars_per_line=14)
            DIV_LINE_H = px(90)
            div_h = DIV_LINE_H * div_lines
            div_y = px(540) - (div_h // 2)
            txt(sl, px(134), div_y, px(900), div_h,
                div_title, 44, bold=True, color=WHITE)
            if sd.get("subtitle"):
                sub_y = div_y + div_h + px(24)
                txt(sl, px(134), sub_y, px(900), px(52),
                    sd["subtitle"], 20, color=WHITE)

        elif stype == "2col":
            sl = new_slide(prs, "Blank")
            grad_bg(sl)
            bot = title_block(sl, sd.get("title", ""), sd.get("subtitle"))
            logo_primary(sl, A)
            two_col(sl, bot,
                    sd.get("left_label", "LEFT"),
                    sd.get("left_body", ""),
                    sd.get("right_label", "RIGHT"),
                    sd.get("right_body", ""),
                    sd.get("footer"))

        elif stype == "3col":
            sl = new_slide(prs, "Blank")
            grad_bg(sl)
            bot = title_block(sl, sd.get("title", ""), sd.get("subtitle"))
            logo_primary(sl, A)
            three_col(sl, bot, sd.get("cols", []), sd.get("footer"))

        elif stype == "content":
            sl = new_slide(prs, "Blank")
            grad_bg(sl)
            bot = title_block(sl, sd.get("title", ""), sd.get("subtitle"))
            logo_primary(sl, A)
            if sd.get("body"):
                txt(sl, px(72), bot + px(28), W - px(144), H - bot - px(80),
                    sd["body"], 14, color=C_GRAY_1)

        elif stype == "closing":
            sl = new_slide(prs, "Blank")
            bg_img(sl, A.get("COVER_3", cover_bg))
            logo_white(sl, A)
            # NO watermark on closing slide
            lines = [
                ("Dankeschön", 44, False),
                ("Thank You", 54, True),
                ("Gracias", 44, False),
                ("\u0927\u0928\u094d\u092f\u0935\u093e\u0926", 40, False),
                ("\u3042\u308a\u304c\u3068\u3046\u3054\u3056\u3044\u307e\u3059", 36, False),
            ]
            cy = px(290)
            for (t, sz, bd) in lines:
                lh = px(int(sz * 1.6))
                txt(sl, 0, cy, W, lh, t, sz, bold=bd, color=WHITE, align=PP_ALIGN.CENTER)
                cy += lh + px(6)
            txt(sl, px(1517), px(957), px(330), px(36), "www.contabo.com", 14, color=WHITE)
            closing_bar(sl)

    # 5. Save
    prs.save(output_path)
    print(f"OK: {output_path} ({len(prs.slides)} slides)")


# ══════════════════════════════════════════════════════════════════════════════
# CLI ENTRY POINT
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python3 build_deck.py content.json output.pptx [master.pptx]")
        sys.exit(1)

    content_path = sys.argv[1]
    output_path = sys.argv[2]
    master_path = sys.argv[3] if len(sys.argv) > 3 else None

    build_deck(content_path, output_path, master_path)
